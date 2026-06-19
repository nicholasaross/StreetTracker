#!/usr/bin/env python3
"""Generate StreetTracker Control Panel documentation deck.

Usage:
    uv run --with python-pptx python scripts/gen_control_deck.py
Output:
    StreetTracker_ControlPanel.pptx  (project root)
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
BG     = RGBColor(0x13, 0x15, 0x1e)
CARD   = RGBColor(0x1c, 0x1f, 0x2e)
CARD2  = RGBColor(0x22, 0x25, 0x38)
BORDER = RGBColor(0x2c, 0x30, 0x4a)
ACCENT = RGBColor(0x4d, 0x9d, 0xff)
GREEN  = RGBColor(0x4c, 0xb5, 0x6b)
AMBER  = RGBColor(0xff, 0xa0, 0x28)
RED    = RGBColor(0xf0, 0x50, 0x50)
PURPLE = RGBColor(0x9b, 0x59, 0xb6)
WHITE  = RGBColor(0xf0, 0xf0, 0xf8)
MUTED  = RGBColor(0x80, 0x84, 0xa0)
CODE_B = RGBColor(0x0c, 0x0e, 0x18)
CODE_F = RGBColor(0x7e, 0xd3, 0x95)

SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
_blank = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def new_slide():
    return prs.slides.add_slide(_blank)

def set_bg(slide, c=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = c

def _no_line(shape):
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn("a:noFill"))

def box(slide, l, t, w, h, fill=None, border=None, bw=Pt(0.75)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        _no_line(s)
    return s

def txt(slide, text, l, t, w, h, size=Pt(14), bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def bullets(slide, items, l, t, w, h, size=Pt(13), color=WHITE):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        lvl, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        prefix = {0: "•", 1: "  ◦", 2: "    ▪"}.get(lvl, "•")
        r.text = f"{prefix}  {text}"
        r.font.size = size
        r.font.color.rgb = color if lvl == 0 else MUTED
    return tb

def code(slide, text, l, t, w, h):
    box(slide, l, t, w, h, fill=CODE_B, border=BORDER, bw=Pt(0.5))
    txt(slide, text, l + Inches(0.18), t + Inches(0.1),
        w - Inches(0.36), h - Inches(0.2),
        size=Pt(10.5), color=CODE_F, wrap=True)

def pill(slide, text, l, t, c=GREEN, w=Inches(0.8)):
    box(slide, l, t, w, Inches(0.28), fill=c)
    txt(slide, text, l + Inches(0.04), t + Inches(0.03),
        w - Inches(0.08), Inches(0.22),
        size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def header(slide, title, sub=None):
    box(slide, 0, 0, SW, Inches(0.82), fill=CARD)
    box(slide, 0, 0, Inches(0.07), Inches(0.82), fill=ACCENT)
    txt(slide, title, Inches(0.22), Inches(0.1), Inches(11), Inches(0.45),
        size=Pt(22), bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.22), Inches(0.52), Inches(11), Inches(0.28),
            size=Pt(11.5), color=MUTED)
    txt(slide, "StreetTracker Control Panel",
        Inches(10.6), Inches(0.28), Inches(2.6), Inches(0.28),
        size=Pt(9.5), color=BORDER, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
box(s, 0, 0, Inches(0.55), SH, fill=ACCENT)
box(s, Inches(0.85), Inches(1.85), Inches(12.0), Inches(1.9), fill=CARD, border=BORDER)
txt(s, "StreetTracker", Inches(1.05), Inches(2.0), Inches(11), Inches(0.72),
    size=Pt(52), bold=True, color=ACCENT)
txt(s, "Control Panel", Inches(1.05), Inches(2.72), Inches(11), Inches(0.72),
    size=Pt(52), bold=True, color=WHITE)
txt(s, "Operator Guide — drive the pull → enrich → train → promote loop",
    Inches(1.05), Inches(3.9), Inches(10), Inches(0.5),
    size=Pt(18), color=MUTED)
txt(s, "Zero tokens consumed once running",
    Inches(1.05), Inches(4.45), Inches(8), Inches(0.4),
    size=Pt(16), italic=True, color=GREEN)
box(s, 0, Inches(6.55), SW, Inches(0.95), fill=CARD)
txt(s, "Dev Box: http://localhost:8095   ·   Showcase: http://127.0.0.1:8090   ·   Orin: http://<orin>:8080",
    Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.42),
    size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# SLIDE 2 — Why It Exists
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Why the Control Panel?",
       "The day-to-day StreetTracker loop, tokenlessly")

box(s, Inches(0.3), Inches(1.1), Inches(6.1), Inches(5.8), fill=CARD, border=BORDER)
txt(s, "Without it", Inches(0.5), Inches(1.22), Inches(5.5), Inches(0.38),
    size=Pt(16), bold=True, color=RED)
bullets(s, [
    "Each session pull requires Claude typing SSH commands",
    "ALPR → DVSA → vehicles → makemodel chain needs Claude attached",
    "Training a new model means watching stdout in a terminal",
    "Promoting a model requires Claude Code on the dev box",
    "Every step burns tokens; operator must wait for each one",
    "No visibility into enrichment state across sessions",
    "No ETA — was that training run hung or just slow?",
], Inches(0.5), Inches(1.7), Inches(5.7), Inches(4.8), size=Pt(13))

box(s, Inches(6.8), Inches(1.1), Inches(6.2), Inches(5.8), fill=CARD, border=BORDER)
txt(s, "With the control panel", Inches(7.0), Inches(1.22), Inches(5.8), Inches(0.38),
    size=Pt(16), bold=True, color=GREEN)
bullets(s, [
    "Launch once → runs detached, survives terminal close",
    "Live radiator: Orin status, sessions, corpus, model",
    "One-click playbooks execute the full enrichment chain",
    "Real progress bars and ETAs — not spinners",
    "Live CNN loss curve while training runs (/training page)",
    "Copy-paste Claude prompt if anything fails",
    "Dashboards viewable from any LAN device (phone, TV)",
    "All control actions are localhost-only (enforced server-side)",
    "Job and playbook history survives a panel restart",
], Inches(7.0), Inches(1.7), Inches(5.8), Inches(4.8), size=Pt(13))

# ---------------------------------------------------------------------------
# SLIDE 3 — Architecture
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Architecture",
       "aiohttp + jinja2 · no new dependencies · mirrors showcase server · per-lane subprocess runner")

# Orin box
box(s, Inches(0.3), Inches(1.1), Inches(3.8), Inches(5.8), fill=CARD, border=ACCENT, bw=Pt(2))
txt(s, "Jetson Orin Nano", Inches(0.5), Inches(1.22), Inches(3.4), Inches(0.38),
    size=Pt(14), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
bullets(s, [
    "streettracker.service (systemd)",
    "RTSP → YOLO+BotSORT → tracks",
    "4K Reolink HTTP snapshots",
    "session_* output files (JSONL)",
    ":8080  live session HTML",
], Inches(0.5), Inches(1.7), Inches(3.4), Inches(3.0), size=Pt(12))
txt(s, "← SSH / scp →", Inches(0.5), Inches(5.2), Inches(3.4), Inches(0.3),
    size=Pt(11), italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Arrow
txt(s, "⇐ SSH / scp ⇒", Inches(4.2), Inches(3.35), Inches(1.5), Inches(0.38),
    size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

# Dev box
box(s, Inches(5.9), Inches(1.1), Inches(7.1), Inches(5.8), fill=CARD, border=BORDER)
txt(s, "Dev Box (Windows · RTX 3080)", Inches(6.1), Inches(1.22), Inches(6.6), Inches(0.38),
    size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

ui_items = [
    (":8080", "Live Dashboard", "served by Orin automatically", MUTED),
    (":8090", "Showcase", "enriched car gallery + /stats", ACCENT),
    (":8095", "Control Panel", "THIS app — operator cockpit", GREEN),
]
for i, (port, name, desc, col) in enumerate(ui_items):
    bx = Inches(6.1) + i * Inches(2.27)
    box(s, bx, Inches(1.75), Inches(2.15), Inches(2.9), fill=CARD2, border=col, bw=Pt(1.5))
    txt(s, port, bx + Inches(0.1), Inches(1.87), Inches(1.9), Inches(0.3),
        size=Pt(14), bold=True, color=col)
    txt(s, name, bx + Inches(0.1), Inches(2.22), Inches(2.0), Inches(0.38),
        size=Pt(13), bold=True, color=WHITE)
    txt(s, desc, bx + Inches(0.1), Inches(2.65), Inches(2.0), Inches(0.38),
        size=Pt(11), color=MUTED)

bullets(s, [
    "Bind: 0.0.0.0 (LAN-viewable dashboards) — but every control POST checks request.remote ∈ loopback",
    "Background pollers: Orin probe ~45 s · local rescan ~30 s · mtime-keyed read cache",
    "GPU lane (alpr / makemodel / build / train) + Net lane (pull / dvsa / vehicles) run concurrently",
    "Windows wake-lock: SetThreadExecutionState(ES_SYSTEM_REQUIRED) while any job runs",
    "On-disk history: control_jobs.json + control_playbooks.json (capped, atomic write)",
], Inches(6.1), Inches(4.8), Inches(6.9), Inches(2.0), size=Pt(12))

# ---------------------------------------------------------------------------
# SLIDE 4 — Quick Start
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Quick Start", "On-demand launcher — no persistent service required")

txt(s, "Recommended: detached launcher (survives terminal close, opens browser automatically)",
    Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.35), size=Pt(14), color=MUTED)
code(s,
     "# Start control panel detached + open http://localhost:8095 in the browser\n"
     "pwsh -NoProfile -File scripts/control_panel.ps1",
     Inches(0.4), Inches(1.48), Inches(12.5), Inches(0.9))

txt(s, "Or launch in the foreground (Ctrl-C to stop):",
    Inches(0.4), Inches(2.55), Inches(12.5), Inches(0.35), size=Pt(14), color=MUTED)
code(s,
     "uv run streettracker control --output-root output\n"
     "# --host 0.0.0.0    (default) LAN-accessible dashboards\n"
     "# --host 127.0.0.1  fully local only\n"
     "# --port 8095        (default)",
     Inches(0.4), Inches(2.98), Inches(12.5), Inches(1.1))

box(s, Inches(0.4), Inches(4.25), Inches(12.5), Inches(3.0), fill=CARD, border=BORDER)
txt(s, "Behaviour on startup:", Inches(0.6), Inches(4.38), Inches(11), Inches(0.35),
    size=Pt(14), bold=True, color=WHITE)
bullets(s, [
    "Panel starts cold — Orin status and local session lists populate within ~45 s of the first background poll",
    "All control actions (submitting jobs, running playbooks) are restricted to localhost — LAN viewers see the radiator but cannot operate",
    "A Windows wake-lock is asserted while any job runs — the dev box will not sleep mid-training",
    "Job and playbook history persists to output/control_jobs.json + control_playbooks.json — survives a panel restart",
    "The showcase site (port 8090) is separate; the 'Re-infer All' playbook can POST a refresh to it when complete",
], Inches(0.6), Inches(4.78), Inches(12.1), Inches(2.35), size=Pt(13))

# ---------------------------------------------------------------------------
# SLIDE 5 — Orin Status Panel
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Dashboard — Orin Status",
       "Live device probe via a single SSH round-trip (~45 s refresh)")

box(s, Inches(0.3), Inches(1.1), Inches(8.2), Inches(5.8), fill=CARD, border=BORDER)
txt(s, "Orin Status", Inches(0.5), Inches(1.22), Inches(7), Inches(0.38),
    size=Pt(16), bold=True, color=WHITE)

pill(s, "● ACTIVE", Inches(0.5), Inches(1.72), GREEN, w=Inches(1.18))
txt(s, "streettracker.service — running  ·  uptime: 6 h 24 m",
    Inches(1.78), Inches(1.75), Inches(6.3), Inches(0.3), size=Pt(13), color=WHITE)

box(s, Inches(0.5), Inches(2.18), Inches(7.5), Inches(0.7), fill=CARD2, border=BORDER, bw=Pt(0.5))
txt(s, "Live session:", Inches(0.65), Inches(2.28), Inches(2.0), Inches(0.28),
    size=Pt(12), color=MUTED)
txt(s, "session_20260616_083047", Inches(2.0), Inches(2.28), Inches(4.5), Inches(0.28),
    size=Pt(12), bold=True, color=ACCENT)
txt(s, "3,847 images banked    VEH: 3,612  ·  PPL: 235",
    Inches(0.65), Inches(2.55), Inches(7.2), Inches(0.28), size=Pt(12), color=MUTED)

box(s, Inches(0.5), Inches(3.0), Inches(7.5), Inches(1.05), fill=CARD2, border=BORDER, bw=Pt(0.5))
txt(s, "Pull estimate (main snaps only):", Inches(0.65), Inches(3.1), Inches(5), Inches(0.28),
    size=Pt(12), color=MUTED)
txt(s, "~42.1 GB  ·  ~7.8 min via local WiFi  ·  3,612 main snaps",
    Inches(0.65), Inches(3.42), Inches(7.2), Inches(0.32), size=Pt(13), bold=True, color=WHITE)
txt(s, "Estimated as of 09:14:22  ·  loaded on demand (not polled in background)",
    Inches(0.65), Inches(3.78), Inches(7.2), Inches(0.22),
    size=Pt(10), color=MUTED, italic=True)

txt(s, "→  Status checked 43 s ago", Inches(0.5), Inches(4.18), Inches(5), Inches(0.25),
    size=Pt(11), italic=True, color=MUTED)

box(s, Inches(0.5), Inches(4.55), Inches(7.5), Inches(1.0), fill=CARD2, border=RED, bw=Pt(0.75))
txt(s, "Unreachable state:", Inches(0.65), Inches(4.65), Inches(4), Inches(0.28),
    size=Pt(12), bold=True, color=RED)
txt(s, "When SSH fails the panel degrades gracefully — last-known status shown with timestamp,\n"
       "no crash, poller retries on the next cycle.",
    Inches(0.65), Inches(4.97), Inches(7.2), Inches(0.48), size=Pt(12), color=MUTED)

bullets(s, [
    "SSH probe combines: systemctl is-active + show -p ActiveEnterTimestamp + session find + VEH/PPL snap counts — one round-trip",
    "Pull estimate calls remote_inventory() over SSH: byte total + file count + main-snap count for --only-main",
    "Estimate is loaded on demand (not polled) — avoids SSH pressure on the live camera device",
    "vehicle_classes = [0,1,2,3,5,7] — people tracked separately from cars/bikes/buses/trucks",
    "service_active = 'active' → green badge  ·  anything else → amber/red with the raw state string",
], Inches(8.7), Inches(1.25), Inches(4.4), Inches(5.5), size=Pt(12))

# ---------------------------------------------------------------------------
# SLIDE 6 — Sessions Table
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Dashboard — Sessions Table",
       "All local sessions, newest first (#N), with enrichment state at a glance")

for col_x, lbl, cw in [
    (Inches(0.3), "#", Inches(0.42)),
    (Inches(0.76), "Session", Inches(3.55)),
    (Inches(4.35), "Duration", Inches(1.0)),
    (Inches(5.4), "Cars", Inches(0.72)),
    (Inches(6.17), "People", Inches(0.72)),
    (Inches(6.95), "Enrichment", Inches(3.6)),
    (Inches(10.6), "Snaps", Inches(0.82)),
]:
    txt(s, lbl, col_x, Inches(1.05), cw, Inches(0.28),
        size=Pt(11), bold=True, color=MUTED)

box(s, Inches(0.3), Inches(1.35), Inches(12.7), Pt(1), fill=BORDER)

rows = [
    ("#17", "session_20260616_083047", "6 h 24 m", "2,841", "188",
     ["ALPR", "DVSA"], "3,612", True),
    ("#16", "session_20260613_092315", "18 h 41 m", "8,203", "512",
     ["ALPR", "DVSA", "VEH", "MM"], "12,841", False),
    ("#15", "session_20260608_125639", "56 h 03 m", "21,417", "1,344",
     ["ALPR", "DVSA", "VEH", "MM"], "33,091", False),
    ("#14", "session_20260601_173815", "72 h 11 m", "19,893", "1,209",
     ["ALPR", "DVSA", "VEH", "MM"], "31,200", False),
    ("#13", "session_20260530_165958", "10 h 02 m", "4,711", "283",
     ["ALPR", "DVSA", "VEH", "MM"], "7,820", False),
]
badge_clr = {"ALPR": ACCENT, "DVSA": GREEN, "VEH": AMBER, "MM": PURPLE}

for i, (num, lbl, dur, cars, ppl, bdgs, snaps, live) in enumerate(rows):
    ry = Inches(1.5) + i * Inches(1.05)
    bg_c = CARD if i % 2 == 0 else CARD2
    box(s, Inches(0.3), ry, Inches(12.7), Inches(0.98), fill=bg_c, border=BORDER, bw=Pt(0.5))
    if live:
        box(s, Inches(0.3), ry, Inches(0.07), Inches(0.98), fill=GREEN)
    txt(s, num, Inches(0.38), ry + Inches(0.35), Inches(0.38), Inches(0.28),
        size=Pt(11), color=GREEN if live else MUTED, bold=live)
    txt(s, lbl, Inches(0.78), ry + Inches(0.35), Inches(3.45), Inches(0.28),
        size=Pt(11), color=WHITE, bold=live)
    txt(s, dur, Inches(4.38), ry + Inches(0.35), Inches(0.92), Inches(0.28),
        size=Pt(11), color=MUTED)
    txt(s, cars, Inches(5.42), ry + Inches(0.35), Inches(0.68), Inches(0.28),
        size=Pt(11), color=WHITE)
    txt(s, ppl, Inches(6.2), ry + Inches(0.35), Inches(0.68), Inches(0.28),
        size=Pt(11), color=MUTED)
    bx = Inches(6.98)
    for b in bdgs:
        pill(s, b, bx, ry + Inches(0.33), badge_clr.get(b, MUTED), w=Inches(0.72))
        bx += Inches(0.8)
    txt(s, snaps, Inches(10.62), ry + Inches(0.35), Inches(0.78), Inches(0.28),
        size=Pt(11), color=MUTED)

box(s, Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.55), fill=CARD, border=BORDER, bw=Pt(0.5))
legend = [("ALPR", ACCENT, "plate reads done"),
          ("DVSA", GREEN, "make/model from MOT API"),
          ("VEH", AMBER, "vehicles aggregated"),
          ("MM", PURPLE, "CNN make/model inferred")]
lx = Inches(0.5)
txt(s, "Badges:", lx, Inches(6.88), Inches(0.75), Inches(0.28),
    size=Pt(11), color=MUTED, bold=True)
lx += Inches(0.88)
for b, c, d in legend:
    pill(s, b, lx, Inches(6.88), c, w=Inches(0.7))
    txt(s, d, lx + Inches(0.78), Inches(6.9), Inches(2.0), Inches(0.26),
        size=Pt(10), color=MUTED)
    lx += Inches(2.82)

# ---------------------------------------------------------------------------
# SLIDE 7 — Corpus & Model
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Dashboard — Corpus & Model",
       "Production model info, corpus size, training run history, and auto recommendations")

# Left: model + corpus
box(s, Inches(0.3), Inches(1.1), Inches(5.9), Inches(5.8), fill=CARD, border=BORDER)
txt(s, "Production Model", Inches(0.5), Inches(1.22), Inches(5.4), Inches(0.38),
    size=Pt(15), bold=True, color=WHITE)
for j, (k, v) in enumerate([
    ("Architecture", "EfficientNet-B5"),
    ("Input size", "456 px"),
    ("Makes", "36"),
    ("make@1 (val)", "0.402"),
    ("Source", "sidecar .meta.json"),
]):
    ry = Inches(1.68) + j * Inches(0.33)
    txt(s, k + ":", Inches(0.5), ry, Inches(1.55), Inches(0.28), size=Pt(12), color=MUTED)
    txt(s, v, Inches(2.15), ry, Inches(3.6), Inches(0.28), size=Pt(12), bold=True, color=WHITE)

box(s, Inches(0.3), Inches(3.42), Inches(5.9), Pt(1), fill=BORDER)
txt(s, "Latest Corpus", Inches(0.5), Inches(3.57), Inches(5.4), Inches(0.38),
    size=Pt(15), bold=True, color=WHITE)
for j, (k, v) in enumerate([
    ("Directory", "runs/uk_crops_0614_512"),
    ("Crops", "32,777"),
    ("Cars (unique)", "~4,100"),
    ("Makes", "36"),
]):
    ry = Inches(4.03) + j * Inches(0.33)
    txt(s, k + ":", Inches(0.5), ry, Inches(1.55), Inches(0.28), size=Pt(12), color=MUTED)
    txt(s, v, Inches(2.15), ry, Inches(3.6), Inches(0.28), size=Pt(12), bold=True, color=WHITE)

# Right: training runs + recommendations
box(s, Inches(6.4), Inches(1.1), Inches(6.6), Inches(2.9), fill=CARD, border=BORDER)
txt(s, "Training Runs", Inches(6.6), Inches(1.22), Inches(6.0), Inches(0.38),
    size=Pt(15), bold=True, color=WHITE)
for j, (rn, bb, acc, col, tag) in enumerate([
    ("uk_make_0614_b5", "B5", "0.402", GREEN, "★ PROD"),
    ("uk_make_0611_b5", "B5", "0.349", MUTED, ""),
    ("uk_make_0608_b0", "B0", "0.271", MUTED, ""),
]):
    ry = Inches(1.68) + j * Inches(0.65)
    box(s, Inches(6.5), ry, Inches(6.3), Inches(0.58), fill=CARD2, border=BORDER, bw=Pt(0.5))
    txt(s, rn, Inches(6.65), ry + Inches(0.06), Inches(3.0), Inches(0.25),
        size=Pt(11), bold=True, color=WHITE)
    txt(s, bb, Inches(9.72), ry + Inches(0.06), Inches(0.6), Inches(0.25),
        size=Pt(11), color=MUTED)
    txt(s, "make@1 " + acc, Inches(10.35), ry + Inches(0.06), Inches(1.55), Inches(0.25),
        size=Pt(12), bold=True, color=col)
    if tag:
        txt(s, tag, Inches(11.9), ry + Inches(0.06), Inches(0.9), Inches(0.25),
            size=Pt(10), color=GREEN)

box(s, Inches(6.4), Inches(4.15), Inches(6.6), Inches(2.75), fill=CARD, border=BORDER)
txt(s, "Recommendations", Inches(6.6), Inches(4.27), Inches(6.0), Inches(0.38),
    size=Pt(15), bold=True, color=WHITE)
bullets(s, [
    "Retrain when corpus grows ≥15 % cars  OR  ≥3 new makes  OR  ≥5 k new crops since the model's training fingerprint",
    "Promote when any candidate run beats production make@1 by ≥0.5 pp",
    "Panel shows the raw evidence numbers, e.g. 'corpus is +18 % cars (4,859 vs 4,100 training) → RETRAIN'",
    "Thresholds are configurable in introspect.py",
], Inches(6.6), Inches(4.72), Inches(6.2), Inches(2.05), size=Pt(12))

# ---------------------------------------------------------------------------
# SLIDE 8 — Individual Jobs
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Process Control — Individual Jobs",
       "Submit, lane, parse stdout live, stream to UI, cancel safely")

# Lane header
box(s, Inches(0.3), Inches(1.08), Inches(12.7), Inches(0.7), fill=CARD, border=BORDER)
txt(s, "GPU Lane (serial):", Inches(0.5), Inches(1.18), Inches(2.2), Inches(0.28),
    size=Pt(12), bold=True, color=ACCENT)
for j, k in enumerate(["alpr-run", "makemodel", "makemodel-build-uk", "makemodel-train-uk", "export-engine"]):
    pill(s, k, Inches(2.8) + j * Inches(1.95), Inches(1.18), ACCENT, w=Inches(1.8))

box(s, Inches(0.3), Inches(1.82), Inches(12.7), Inches(0.7), fill=CARD2, border=BORDER)
txt(s, "Net Lane (serial):", Inches(0.5), Inches(1.92), Inches(2.2), Inches(0.28),
    size=Pt(12), bold=True, color=GREEN)
for j, k in enumerate(["pull", "dvsa-label", "dvsa-apply", "vehicles"]):
    pill(s, k, Inches(2.8) + j * Inches(1.95), Inches(1.92), GREEN, w=Inches(1.8))

txt(s, "GPU + Net lanes run concurrently — a pull overlaps a train; two GPU jobs never do",
    Inches(0.4), Inches(2.6), Inches(12), Inches(0.28),
    size=Pt(11), italic=True, color=MUTED)

# Mock job card
box(s, Inches(0.3), Inches(3.0), Inches(8.0), Inches(4.3), fill=CARD, border=BORDER)
txt(s, "Live job card:", Inches(0.5), Inches(3.1), Inches(5), Inches(0.32),
    size=Pt(14), bold=True, color=WHITE)

box(s, Inches(0.45), Inches(3.5), Inches(7.7), Inches(3.55), fill=CARD2, border=BORDER, bw=Pt(0.5))
txt(s, "makemodel-train-uk", Inches(0.62), Inches(3.6), Inches(4.5), Inches(0.32),
    size=Pt(13), bold=True, color=WHITE)
pill(s, "● running", Inches(5.1), Inches(3.58), GREEN, w=Inches(1.15))
txt(s, "epoch 14/30  ·  make@1 0.391  ·  loss 0.42  ·  ETA ~7 m",
    Inches(0.62), Inches(3.96), Inches(7.2), Inches(0.28), size=Pt(12), color=MUTED)
box(s, Inches(0.62), Inches(4.32), Inches(7.18), Inches(0.2), fill=BORDER)
box(s, Inches(0.62), Inches(4.32), Inches(7.18 * 14 / 30), Inches(0.2), fill=ACCENT)
txt(s, "14/30 epochs (46.7 %)", Inches(0.62), Inches(4.58), Inches(4), Inches(0.25),
    size=Pt(11), color=MUTED)
code(s,
     "[makemodel-uk] epoch 14/30 loss=0.4201 make@1=0.391 make@5=0.689 (43s)\n"
     "[makemodel-uk] epoch 13/30 loss=0.4389 make@1=0.383 make@5=0.681 (42s)",
     Inches(0.62), Inches(4.88), Inches(6.85), Inches(0.82))
txt(s, "[ Cancel ]   [ ↗ Full training page ]",
    Inches(0.62), Inches(5.82), Inches(4), Inches(0.28),
    size=Pt(12), color=ACCENT)

bullets(s, [
    "stdout piped live via PYTHONUNBUFFERED=1 child env",
    "Parser: pure text-in → Progress dataclass out (no clock/IO inside parser)",
    "500-line ring buffer → log viewer on demand",
    "Pull: dir-watcher polls local session dir size vs remote byte total → real byte/s + ETA bar",
    "Cancel: taskkill /F /T /PID on Windows (kills uv run + child python process tree)",
    "Wake-lock: SetThreadExecutionState held while ≥1 job runs",
    "Exit 0 with traceback / OOM / config-error in output still surfaces 'Copy Claude prompt'",
], Inches(8.5), Inches(3.0), Inches(4.6), Inches(4.2), size=Pt(12))

# ---------------------------------------------------------------------------
# SLIDE 9 — Playbooks Overview
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Process Control — Playbooks",
       "Multi-step procedures: sequential, stop-on-failure, rest → skipped")

rows_pb = [
    ("Enrich Session",  ACCENT, "alpr-run  →  dvsa-label  →  dvsa-apply  →  vehicles  →  makemodel",  "Per-session · safe"),
    ("Build + Train",   ACCENT, "makemodel-build-uk  →  makemodel-train-uk  (dated output dirs)",       "Safe · use /training for live curve"),
    ("Roll Session",    AMBER,  "Restart Orin service  →  verify handover + track count  →  pull closed session", "⚠  Confirm required"),
    ("Promote Model",   RED,    "Backup makemodel_b0.pt  →  swap in best.pt  →  write meta sidecar",   "⚠  Confirm required"),
    ("Re-infer All",    GREEN,  "makemodel on every local session  →  POST /api/refresh to showcase",   "Safe · non-destructive"),
]
for i, (name, col, desc, tag) in enumerate(rows_pb):
    ry = Inches(1.12) + i * Inches(1.12)
    box(s, Inches(0.3), ry, Inches(12.7), Inches(1.02), fill=CARD, border=col, bw=Pt(1.8))
    txt(s, name, Inches(0.5), ry + Inches(0.1), Inches(2.3), Inches(0.38),
        size=Pt(16), bold=True, color=col)
    txt(s, tag, Inches(2.95), ry + Inches(0.12), Inches(3.5), Inches(0.32),
        size=Pt(11), italic=True, color=MUTED)
    txt(s, desc, Inches(0.5), ry + Inches(0.55), Inches(12.2), Inches(0.38),
        size=Pt(13), color=WHITE)

txt(s, "Engine: ordered steps run sequentially · first failure stops and marks remaining as skipped · each step is a job (live progress) or an action (SSH/file op)",
    Inches(0.3), Inches(6.78), Inches(12.7), Inches(0.35),
    size=Pt(11), italic=True, color=MUTED)

# ---------------------------------------------------------------------------
# SLIDE 10 — Playbook Detail
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Playbooks — Step Detail",
       "What each step does, what can fail, and what the confirm gate protects")

box(s, Inches(0.3), Inches(1.1), Inches(6.15), Inches(2.65), fill=CARD, border=ACCENT, bw=Pt(1.5))
txt(s, "Enrich Session", Inches(0.5), Inches(1.22), Inches(4), Inches(0.35),
    size=Pt(14), bold=True, color=ACCENT)
for j, step in enumerate([
    "1.  alpr-run  — plate inference with ghost mask + bbox hints",
    "2.  dvsa-label  — DVSA MOT API lookup per canonical plate",
    "3.  dvsa-apply  — write make/model onto per-track records",
    "4.  vehicles  — per-vehicle aggregation (_vehicles.json)",
    "5.  makemodel  — CNN make inference (_makemodel_by_track.json)",
]):
    txt(s, step, Inches(0.5), Inches(1.62) + j * Inches(0.37),
        Inches(5.8), Inches(0.32), size=Pt(12), color=WHITE)

box(s, Inches(6.5), Inches(1.1), Inches(6.55), Inches(2.65), fill=CARD, border=AMBER, bw=Pt(1.5))
txt(s, "Roll Session  ⚠ Confirm", Inches(6.7), Inches(1.22), Inches(6.1), Inches(0.35),
    size=Pt(14), bold=True, color=AMBER)
bullets(s, [
    "Action: probe Orin — refuse if not actively recording",
    "Action: sudo systemctl restart streettracker.service",
    "Action: verify new session label ≠ old, count finalised tracks",
    "Job: pull --session <old_session> --only-main --skip-existing",
], Inches(6.7), Inches(1.62), Inches(6.15), Inches(2.0), size=Pt(12))

box(s, Inches(0.3), Inches(3.9), Inches(6.15), Inches(2.55), fill=CARD, border=ACCENT, bw=Pt(1.5))
txt(s, "Build + Train", Inches(0.5), Inches(4.02), Inches(4), Inches(0.35),
    size=Pt(14), bold=True, color=ACCENT)
bullets(s, [
    "Job: makemodel-build-uk runs/uk_crops_<date>_512  — discovers all local sessions, harvests DVSA-labelled crops into a manifest",
    "Job: makemodel-train-uk <corpus> --backbone b5 --input-size 456  — live loss curve visible on /training",
    "Output: runs/uk_make_<date>_b5/best.pt  +  history.json",
], Inches(0.5), Inches(4.45), Inches(5.8), Inches(1.85), size=Pt(12))

box(s, Inches(6.5), Inches(3.9), Inches(6.55), Inches(2.55), fill=CARD, border=RED, bw=Pt(1.5))
txt(s, "Promote Model  ⚠ Confirm", Inches(6.7), Inches(4.02), Inches(6.1), Inches(0.35),
    size=Pt(14), bold=True, color=RED)
bullets(s, [
    "Action: read best run's history.json (best_epoch, val make@1, makes list)",
    "Action: backup makemodel_b0.pt → makemodel_b0.<date>_<acc>.pt",
    "Action: copy best.pt → makemodel_b0.pt",
    "Action: write makemodel_b0.meta.json sidecar (arch, input_size, n_makes, val_make_top1, source_run)",
], Inches(6.7), Inches(4.45), Inches(6.15), Inches(1.85), size=Pt(12))

txt(s, "Confirm gate: server refuses roll / promote without confirm=true in POST body  ·  UI shows window.confirm() dialog  ·  LAN viewers see no submit button",
    Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.5),
    size=Pt(11), italic=True, color=MUTED)

# ---------------------------------------------------------------------------
# SLIDE 11 — Training Page
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "/training — Full-Page Training View",
       "Live CNN loss curve, epoch table, corpus distribution, one-click promote")

box(s, Inches(0.3), Inches(1.1), Inches(8.0), Inches(5.95), fill=CARD, border=BORDER)

# Pre-training facts
txt(s, "Pre-training facts", Inches(0.5), Inches(1.22), Inches(5), Inches(0.32),
    size=Pt(13), bold=True, color=MUTED)
for j, (k, v) in enumerate([("Backbone", "EfficientNet-B5"), ("Input", "456 px"),
                              ("Makes", "36"), ("Train", "26,500 crops"),
                              ("Val", "6,277 crops"), ("Epochs", "30")]):
    x = Inches(0.5) + (j % 3) * Inches(2.48)
    y = Inches(1.6) + (j // 3) * Inches(0.35)
    txt(s, k + ":", x, y, Inches(1.15), Inches(0.28), size=Pt(11), color=MUTED)
    txt(s, v, x + Inches(1.2), y, Inches(1.2), Inches(0.28), size=Pt(11), bold=True, color=WHITE)

# Chart area
box(s, Inches(0.45), Inches(2.42), Inches(7.5), Inches(2.95), fill=CODE_B, border=BORDER, bw=Pt(0.5))
txt(s, "Dual axis: loss (amber) + make@1 (blue) per epoch — updates live each epoch",
    Inches(0.62), Inches(2.52), Inches(7.0), Inches(0.28), size=Pt(11), color=MUTED)
txt(s, "loss  ▼", Inches(0.62), Inches(2.9), Inches(1.5), Inches(0.28),
    size=Pt(13), color=AMBER)
txt(s, "make@1  ▲", Inches(0.62), Inches(3.35), Inches(2.0), Inches(0.28),
    size=Pt(13), color=ACCENT)
txt(s, "- - - - - - production make@1 = 0.402  (dashed reference line)",
    Inches(0.62), Inches(3.8), Inches(7.0), Inches(0.28),
    size=Pt(11), color=MUTED, italic=True)
txt(s, "Epoch 14/30  ·  best make@1: 0.391 @ epoch 12  ·  ETA ~7 m",
    Inches(0.62), Inches(5.18), Inches(7.0), Inches(0.28), size=Pt(12), color=WHITE)

# Promote button (completion state)
box(s, Inches(0.45), Inches(5.55), Inches(7.5), Inches(0.62), fill=CARD2, border=GREEN, bw=Pt(1.0))
txt(s, "On completion:  [ Promote this run ]  (localhost, confirm-gated) appears alongside candidate vs production verdict",
    Inches(0.62), Inches(5.68), Inches(7.2), Inches(0.38), size=Pt(12), color=GREEN)

# Right
box(s, Inches(8.5), Inches(1.1), Inches(4.8), Inches(5.95), fill=CARD, border=BORDER)
txt(s, "How to reach it", Inches(8.7), Inches(1.22), Inches(4.3), Inches(0.32),
    size=Pt(14), bold=True, color=WHITE)
code(s, "http://localhost:8095/training",
     Inches(8.6), Inches(1.62), Inches(4.55), Inches(0.52))
bullets(s, [
    "Shows the current or most-recent makemodel-train-uk job",
    "Chart renders from parsed per-epoch Progress metrics in-memory — no extra file I/O during training",
    "Dashed reference line = production model make@1 read from the sidecar",
    "Corpus distribution: per-make crop bar chart with brand SVG logos",
    "Per-epoch table: epoch, loss, make@1, make@5, elapsed sec",
    "Nav bar links Dashboard ↔ Training",
    "Polls /api/jobs every ~5 s; training chart auto-extends as new epochs arrive",
    "Build-Train playbook runs this view automatically for the running job",
], Inches(8.7), Inches(2.25), Inches(4.45), Inches(4.65), size=Pt(12))

# ---------------------------------------------------------------------------
# SLIDE 12 — Failure Recovery & Security
# ---------------------------------------------------------------------------
s = new_slide(); set_bg(s)
header(s, "Failure Recovery & Security Model",
       "Copy-paste Claude prompt on failure · localhost-only control · LAN dashboards")

box(s, Inches(0.3), Inches(1.1), Inches(6.35), Inches(5.95), fill=CARD, border=RED, bw=Pt(1.5))
txt(s, "When a Job Fails", Inches(0.5), Inches(1.22), Inches(5.8), Inches(0.38),
    size=Pt(16), bold=True, color=RED)
bullets(s, [
    "Each failed (or suspicious-exit-0) job shows a 'Copy Claude prompt' button",
    "One click → clipboard has a complete, paste-ready help request",
    "Prompt content: project path, exact command, exit code, elapsed time, captured output tail, kind-specific diagnosis pointer",
    "Triggered by: non-zero exit · traceback · CUDA OOM · config-error in output",
    "Prompt survives a panel restart (persisted in control_jobs.json)",
], Inches(0.5), Inches(1.68), Inches(5.95), Inches(2.25), size=Pt(13))

code(s,
     "Project: StreetTracker  (D:\\Projects\\StreetTracker)\n"
     "Command: uv run streettracker pull --session session_x --only-main\n"
     "Status:  FAILED (exit 1)  |  Elapsed: 73 s\n"
     "Issue:   could not reach the device over SSH\n"
     "\n"
     "Output (last 20 lines):\n"
     "```\n"
     "[pull] session: session_x\n"
     "ssh: connect to host orin port 22: timed out\n"
     "```\n"
     "\n"
     "See CLAUDE.md → Ops for SSH recovery steps.\n"
     "Diagnose, fix, then re-run to confirm.",
     Inches(0.45), Inches(3.97), Inches(6.05), Inches(2.75))

box(s, Inches(6.85), Inches(1.1), Inches(6.3), Inches(5.95), fill=CARD, border=AMBER, bw=Pt(1.5))
txt(s, "Security Model", Inches(7.05), Inches(1.22), Inches(5.8), Inches(0.38),
    size=Pt(16), bold=True, color=AMBER)
bullets(s, [
    "Server binds 0.0.0.0 — dashboards readable from any LAN device",
    "All mutating routes check request.remote ∈ {127.0.0.1, ::1}",
    "LAN IPs always 403 — regardless of any headers",
    "Destructive playbooks (roll, promote) also require confirm=true in the POST body",
    "UI hides control buttons when is_local=false (server enforces regardless)",
    "--host 127.0.0.1 flag for fully local-only mode",
], Inches(7.05), Inches(1.68), Inches(6.0), Inches(2.7), size=Pt(13))

txt(s, "Persistence:", Inches(7.05), Inches(4.45), Inches(5.5), Inches(0.32),
    size=Pt(14), bold=True, color=WHITE)
bullets(s, [
    "Finished jobs → output/control_jobs.json  (capped at 200, atomic write via tempfile + os.replace)",
    "Finished playbooks → output/control_playbooks.json  (same)",
    "Panel restart: snapshots() merges prior-run history with live items — panel feels continuous",
    "Terminal snapshot of each failed job (incl. copy-prompt text) is preserved in the history file",
], Inches(7.05), Inches(4.85), Inches(6.0), Inches(2.1), size=Pt(13))

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = Path(__file__).resolve().parent.parent / "StreetTracker_ControlPanel.pptx"
prs.save(str(out))
print(f"Saved: {out}")
